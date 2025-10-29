#!/usr/bin/env python3
"""
Script to generate PowerPoint presentation for Chess AI project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title):
    """Add a content slide with title"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    return slide

def add_bullet_points(text_frame, bullets, level=0):
    """Add bullet points to a text frame"""
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = bullet
        p.level = level
        p.font.size = Pt(18) if level == 0 else Pt(16)

def add_image_to_slide(slide, image_path, left, top, width, height):
    """Add image to slide if it exists"""
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)
        return True
    return False

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = add_title_slide(
        prs,
        "Teaching Computers to Play Chess Like Masters",
        "Building an AI Chess Engine with Deep Learning & Classical Search\n\n♟️ ♞ ♝ ♜ ♛ ♚"
    )

    # Slide 2: The Problem
    slide = add_content_slide(prs, "The Problem: Can We Build Smart Chess AI Affordably?")
    content = slide.placeholders[1].text_frame
    add_bullet_points(content, [
        "The Challenge:",
        "Chess engines like AlphaZero are superhuman but require weeks of training",
        "Thousands of specialized processors costing over $100,000",
        "",
        "The Question: Can we create a competitive chess AI that:",
        "Trains in under an hour on a single GPU?",
        "Reaches intermediate club-level strength (~2000 Elo)?",
        "Costs less than $10 to train?",
        "Learns from expert human games instead of expensive self-play?",
        "",
        "Answer: YES! By using supervised learning from master games."
    ])

    # Slide 3: Why Chess?
    slide = add_content_slide(prs, "Why Chess?")
    content = slide.placeholders[1].text_frame
    add_bullet_points(content, [
        "Perfect AI Testbed:",
        "Clear outcomes: Win/loss/draw - perfect feedback",
        "Huge complexity: ~10^43 legal positions",
        "Pattern recognition: Requires understanding, not just calculation",
        "Strong benchmarks: Stockfish, Maia, human Elo ratings",
        "",
        "Rich Data Available:",
        "Lichess Database: 3.1 GB of master games",
        "Millions of positions from expert play",
        "Free & open: CC0 license",
        "Labeled by masters: Each move is a training example"
    ])

    # Add stats
    left = Inches(1)
    top = Inches(5)
    width = Inches(8)
    height = Inches(1)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "10^43 Positions  •  3.7M Training Examples  •  ~35 Avg Legal Moves"
    p.font.size = Pt(20)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Slide 4: Our Approach
    slide = add_content_slide(prs, "Our Approach: Supervised Learning")
    content = slide.placeholders[1].text_frame
    add_bullet_points(content, [
        "Two Paths to Chess AI:",
        "",
        "Our Approach (Supervised):",
        "Learn from expert human games",
        "Training time: 40 minutes on single GPU",
        "Cost: ~$10",
        "Final strength: ~2000 Elo (intermediate)",
        "Simple & reproducible",
        "",
        "AlphaZero (Reinforcement):",
        "Self-play against itself",
        "Training time: Weeks on thousands of TPUs",
        "Cost: ~$100,000+",
        "Final strength: 3000+ Elo (superhuman)",
        "Research-level complexity"
    ])

    # Slide 5: The Data
    slide = add_content_slide(prs, "The Data: Learning from Masters")
    content = slide.placeholders[1].text_frame
    add_bullet_points(content, [
        "Data Collection:",
        "Source: Lichess Elite Database",
        "Size: 3.1 GB of game files",
        "Filter: Both players 2000+ Elo",
        "Result: 3.7 million positions",
        "",
        "Data Processing:",
        "Skip first 8 moves (opening theory)",
        "Balance white/black 50/50",
        "Tag by phase (opening/middle/endgame)",
        "Split: 70% train / 15% validation / 15% test"
    ])

    # Add phase distribution image
    img_path = "reports/figures/phase_distribution.png"
    add_image_to_slide(slide, img_path, Inches(5.5), Inches(2), Inches(4), Inches(3.5))

    # Slide 6: The Neural Network
    slide = add_content_slide(prs, "The Brain: Neural Network Architecture")
    content = slide.placeholders[1].text_frame
    add_bullet_points(content, [
        "Mini-ResNet Design (Inspired by AlphaZero):",
        "",
        "Input: 12×8×8 board tensor (12 piece types)",
        "Conv2D layer (64 channels)",
        "6× Residual Blocks for pattern learning",
        "Split into two heads:",
        "",
        "Policy Head: Predicts move probabilities (4,672 possibilities)",
        "Value Head: Evaluates position (-1 to +1)",
        "",
        "Key Features:",
        "6M parameters, 23MB model size",
        "Residual connections enable deeper learning",
        "Spatial convolutions understand board patterns",
        "Legal move masking - only consider valid moves"
    ])

    # Slide 7: Training Setup
    slide = add_content_slide(prs, "Training: Teaching the Network")
    content = slide.placeholders[1].text_frame
    add_bullet_points(content, [
        "Training Configuration:",
        "Optimizer: AdamW (adaptive learning rate)",
        "Learning rate: Cosine schedule (0.001 → 0.0001)",
        "Batch size: 256 positions",
        "Epochs: 10 passes through data",
        "Hardware: Google Colab L4 GPU",
        "Time: 40 minutes total",
        "",
        "Loss Function:",
        "Policy Loss: How accurate are move predictions?",
        "Value Loss: How accurate is position evaluation?",
        "Total = CrossEntropy(policy) + 0.35 × MSE(value)",
        "",
        "Result: 28.8% top-1 accuracy (10× better than random!)"
    ])

    # Add training curves
    img_path = "reports/figures/training_curves.png"
    add_image_to_slide(slide, img_path, Inches(5.5), Inches(2), Inches(4), Inches(3.5))

    # Slide 8: Training Results
    slide = add_content_slide(prs, "Training Results: Strong Pattern Recognition")
    content = slide.placeholders[1].text_frame
    add_bullet_points(content, [
        "Final Metrics:",
        "Top-1 Accuracy: 28.8% (network's 1st choice matches master)",
        "Top-3 Accuracy: 52% (master's move in top 3 predictions)",
        "Top-5 Accuracy: 64%",
        "",
        "Context:",
        "Random guessing: ~2.9% (with ~35 legal moves)",
        "Human masters (~2200 Elo): estimated 40-50%",
        "",
        "Key Insights:",
        "Data Impact: 37× more data → +44% accuracy improvement",
        "Training Time: Only 40 minutes for intermediate understanding",
        "Efficient: Single GPU achieved strong results",
        "",
        "Lesson: Data quality matters more than architecture tweaks!"
    ])

    # Slide 9: The Search Algorithm
    slide = add_content_slide(prs, "The Search: Combining Neural + Classical AI")
    content = slide.placeholders[1].text_frame
    add_bullet_points(content, [
        "Why Search?",
        "Neural network provides intuition, but chess requires looking ahead",
        "Solution: Alpha-Beta search enhanced with neural guidance",
        "",
        "How It Works:",
        "1. Generate all legal moves",
        "2. Neural ordering: Sort by policy network (best first)",
        "3. Search tree: Explore best moves deeply, prune bad branches",
        "4. Leaf evaluation: Use value network to judge positions",
        "5. Back up scores through the tree",
        "6. Choose the move with best score",
        "",
        "Optimizations:",
        "Transposition table: Remember 100K positions (20-30% hit rate)",
        "Killer moves: Try moves that worked before",
        "Quiescence search: Extend on captures/checks",
        "Search depth: 2-3 plies in 300ms per move"
    ])

    # Slide 10: Benchmark Results
    slide = add_content_slide(prs, "The Results: How Strong Is Our AI?")
    content = slide.placeholders[1].text_frame
    add_bullet_points(content, [
        "Tournament Format: 100 games each, 300ms/move",
        "",
        "vs Sunfish (Depth 2, ~2000 Elo):",
        "Result: 45 wins - 36 draws - 19 losses",
        "Win rate: 63% (+92 Elo)",
        "Strong positional play with high draw rate",
        "",
        "vs Stockfish Level 1 (~2300 Elo):",
        "Result: 12-6-82",
        "Win rate: 15% (-301 Elo)",
        "Expected deficit from deeper search (depth 10-15 vs our 2-3)",
        "",
        "vs Maia-1500 (Neural engine):",
        "Result: 1-2-97",
        "Win rate: 2% (-676 Elo)",
        "",
        "Estimated Elo: ~2000 (Intermediate club level)"
    ])

    # Slide 11: Performance Visualizations
    slide = add_content_slide(prs, "Performance Visualizations")

    # Add Elo estimates image
    img_path = "reports/figures/elo_estimates.png"
    if add_image_to_slide(slide, img_path, Inches(0.5), Inches(1.5), Inches(4.5), Inches(3)):
        # Add ACPL image
        img_path2 = "reports/figures/acpl_by_phase.png"
        add_image_to_slide(slide, img_path2, Inches(5), Inches(1.5), Inches(4.5), Inches(3))

    # Add caption
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "Left: Relative strength vs opponents  •  Right: Avg Centipawn Loss by phase (lower is better)"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER

    # Slide 12: What Did It Learn?
    slide = add_content_slide(prs, "What Did The AI Actually Learn?")
    content = slide.placeholders[1].text_frame
    add_bullet_points(content, [
        "Strengths:",
        "Positional understanding: Controls center, develops pieces",
        "Basic tactics: Recognizes forks, pins, discovered attacks",
        "Opening principles: Strong first 10-15 moves",
        "Pattern recognition: Sees common motifs from masters",
        "",
        "Weaknesses:",
        "Deep tactics (42% of errors): Misses 4+ move combinations",
        "Endgame precision (28%): Suboptimal technical play",
        "Positional nuances (18%): Misses prophylactic moves",
        "King safety (12%): Undervalues defensive moves",
        "",
        "Playing Style:",
        "Like an intermediate club player (Class A/Expert level)",
        "Strong fundamentals but occasionally misses complex tactics"
    ])

    # Slide 13: Key Insights
    slide = add_content_slide(prs, "Key Insights & Takeaways")
    content = slide.placeholders[1].text_frame
    add_bullet_points(content, [
        "1. Supervised Learning Is Practical",
        "40 minutes + $10 achieved intermediate chess strength",
        "You don't need massive resources to build useful AI!",
        "",
        "2. Data Quality Matters Most",
        "37× more data → +44% accuracy improvement",
        "Focus on data before architecture tweaks",
        "",
        "3. Hybrid Approaches Work Best",
        "Neural networks + Classical algorithms = better than either alone",
        "",
        "4. Trade-offs Are Real",
        "Supervised: Fast, cheap, simple → ~2000 Elo",
        "Reinforcement: Slow, expensive, complex → 3000+ Elo",
        "",
        "5. AI Is Increasingly Accessible",
        "What required institutions in 2010 now takes one person and cloud GPUs"
    ])

    # Slide 14: Future Improvements
    slide = add_content_slide(prs, "Future Improvements")
    content = slide.placeholders[1].text_frame
    add_bullet_points(content, [
        "Quick Wins (High ROI):",
        "Opening book: +50-100 Elo (2 days effort)",
        "Endgame tablebases: +30-50 Elo (1 day)",
        "Model quantization: 2-3× faster inference → deeper search",
        "",
        "Medium-Term Improvements:",
        "Scale to 10M positions: +100-150 Elo (2 weeks)",
        "Larger model (20M params): +80-120 Elo (1 week)",
        "Parallel search: +50-100 Elo (1 week)",
        "",
        "Potential with all improvements:",
        "~2300-2400 Elo (Expert/Master level)",
        "",
        "Long-Term (Research):",
        "Self-play reinforcement learning: +200-400 Elo (1-2 months)",
        "NNUE architecture: 10× faster inference"
    ])

    # Slide 15: System Architecture
    slide = add_content_slide(prs, "Complete System Architecture")

    # Add architecture as text
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    tf = textbox.text_frame
    tf.word_wrap = True

    architecture_text = """TRAINING PHASE:
Lichess PGN Files (3.1GB) → Filter & Parse → 3.7M Positions
→ DataLoader → MiniResNet Neural Network (6M params)
→ Loss Function → AdamW Optimizer (40 min) → best_model.pth

INFERENCE PHASE:
Current Position → Encode Board (12×8×8 tensor)
→ Neural Network → Policy logits + Value
→ Alpha-Beta Search (300ms, depth 2-3)
   • Neural move ordering
   • Transposition table (100K entries)
   • Quiescence search on captures
   • Alpha-beta pruning
→ Best Move → Apply to board → Repeat

BENCHMARKING:
Model vs Engines → 100 games each → Statistics → Elo estimates"""

    p = tf.paragraphs[0]
    p.text = architecture_text
    p.font.size = Pt(14)
    p.font.name = 'Courier New'

    # Slide 16: Project Statistics
    slide = add_content_slide(prs, "Project Statistics")

    stats_text = """8,781 Lines of Code  •  33 Python Modules  •  5 Jupyter Notebooks

3.7M Training Positions  •  6M Model Parameters  •  300 Benchmark Games

40 minutes Training  •  ~$10 Total Cost  •  ~2000 Elo Rating


A Complete Deep Learning System
From raw data to trained model to competitive gameplay"""

    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = stats_text
    p.font.size = Pt(22)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Slide 17: Conclusion
    slide = add_content_slide(prs, "Conclusion")
    content = slide.placeholders[1].text_frame
    add_bullet_points(content, [
        "We Successfully Built a Chess AI!",
        "",
        "Problem: Build competitive chess AI affordably",
        "Approach: Supervised learning from 3.7M master positions",
        "Method: Mini-ResNet + Alpha-Beta Search",
        "Result: ~2000 Elo in 40 minutes for $10",
        "",
        "What We Demonstrated:",
        "Deep learning fundamentals",
        "Full ML pipeline (data → training → evaluation)",
        "Hybrid AI systems (neural + classical)",
        "Practical resource constraints",
        "",
        "Broader Impact:",
        "AI doesn't require massive budgets",
        "Supervised learning is practical for real applications",
        "Trade-offs exist between different approaches",
        "Data quality matters most"
    ])

    # Slide 18: Thank You
    slide = add_title_slide(
        prs,
        "Thank You!",
        "Questions?\n\n♟️ 3.7M training positions  •  🧠 6M parameters\n⚡ 40-minute training  •  🏆 ~2000 Elo  •  💰 $10 cost"
    )

    # Save presentation
    output_path = "Chess_AI_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
